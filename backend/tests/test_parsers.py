"""
Parametrized tests for the Phase 5 document parsing layer.

These tests create *real* files on disk (inside a temporary directory)
and feed them through the actual parser classes.  This ensures we are
testing the genuine parsing logic — not mocked stubs.
"""
import pytest
import io
from pathlib import Path
from unittest.mock import patch, MagicMock

# Force SQLAlchemy mapper registration
import authentication.models  # noqa: F401

from parsers.base import Document, BaseParser
from parsers.text_parser import TextParser
from parsers.pdf_parser import PdfParser
from parsers.docx_parser import DocxParser
from parsers.pptx_parser import PptxParser
from parsers.image_parser import ImageParser, CaptionProvider
from parsers.factory import get_parser


# ------------------------------------------------------------------
# Fixtures
# ------------------------------------------------------------------

@pytest.fixture
def tmp_workspace(tmp_path):
    """Return a temporary directory to act as a mini workspace."""
    return tmp_path


# ------------------------------------------------------------------
# BaseParser interface contract
# ------------------------------------------------------------------

def test_base_parser_is_abstract():
    """BaseParser cannot be instantiated directly."""
    with pytest.raises(TypeError):
        BaseParser()


# ------------------------------------------------------------------
# Document dataclass
# ------------------------------------------------------------------

def test_document_defaults():
    doc = Document(filename="a.txt", source_path="a.txt", text="hello")
    assert doc.page_count == 1
    assert doc.metadata == {}


# ------------------------------------------------------------------
# TextParser
# ------------------------------------------------------------------

@pytest.mark.parametrize("ext, content, expected_substring", [
    (".txt", "Hello world", "Hello world"),
    (".md", "# Heading\nSome content", "# Heading"),
    (".markdown", "**bold text**", "**bold text**"),
])
def test_text_parser_valid(tmp_workspace, ext, content, expected_substring):
    filepath = tmp_workspace / f"sample{ext}"
    filepath.write_text(content, encoding="utf-8")

    parser = TextParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path=f"sample{ext}")

    assert isinstance(doc, Document)
    assert expected_substring in doc.text
    assert doc.filename == f"sample{ext}"
    assert doc.metadata["parser"] == "TextParser"


def test_text_parser_empty_file(tmp_workspace):
    filepath = tmp_workspace / "empty.txt"
    filepath.write_text("", encoding="utf-8")

    parser = TextParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="empty.txt")
    assert doc.text == ""


def test_text_parser_latin1_fallback(tmp_workspace):
    """File with bytes invalid in UTF-8 should fall back to latin-1."""
    filepath = tmp_workspace / "latin.txt"
    filepath.write_bytes(b"\xe9\xe8\xea")  # é è ê in latin-1

    parser = TextParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="latin.txt")
    assert len(doc.text) == 3  # Three decoded characters


# ------------------------------------------------------------------
# PdfParser
# ------------------------------------------------------------------

def _create_minimal_pdf(filepath: Path, text: str = "Hello PDF"):
    """Create a tiny valid PDF using pypdf."""
    from pypdf import PdfWriter
    from pypdf._page import PageObject
    from pypdf.generic import (
        ArrayObject,
        DictionaryObject,
        NameObject,
        NumberObject,
        TextStringObject,
        StreamObject,
    )

    writer = PdfWriter()

    # Build a minimal page with a text stream
    content = f"BT /F1 12 Tf 100 700 Td ({text}) Tj ET"
    stream = StreamObject()
    stream.set_data(content.encode("latin-1"))

    font_dict = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    resources = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): font_dict}
            )
        }
    )
    page = PageObject.create_blank_page(width=612, height=792)
    page[NameObject("/Resources")] = resources
    page[NameObject("/Contents")] = stream

    writer.add_page(page)
    writer._add_object(stream)

    with open(filepath, "wb") as f:
        writer.write(f)


@patch("pdf2image.convert_from_bytes")
@patch("parsers.pdf_parser.pytesseract.image_to_string")
def test_pdf_parser_valid(mock_image_to_string, mock_convert, tmp_workspace):
    from PIL import Image
    # Mock returning one blank image page
    mock_convert.return_value = [Image.new("RGB", (100, 100), color="white")]
    mock_image_to_string.return_value = "Hello PDF"
    
    filepath = tmp_workspace / "sample.pdf"
    _create_minimal_pdf(filepath, "Hello PDF")

    parser = PdfParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="sample.pdf")

    assert isinstance(doc, Document)
    assert doc.page_count >= 1
    assert doc.metadata["parser"] == "PdfParser"


def test_pdf_parser_corrupt(tmp_workspace):
    filepath = tmp_workspace / "corrupt.pdf"
    filepath.write_bytes(b"this is not a pdf")

    parser = PdfParser()
    with pytest.raises(Exception):
        parser.parse(filepath.read_bytes(), filename=filepath.name)


# ------------------------------------------------------------------
# DocxParser
# ------------------------------------------------------------------

def _create_minimal_docx(filepath: Path, text: str = "Hello DOCX"):
    from docx import Document as DocxDocument
    doc = DocxDocument()
    doc.add_paragraph(text)
    doc.save(str(filepath))


def test_docx_parser_valid(tmp_workspace):
    filepath = tmp_workspace / "sample.docx"
    _create_minimal_docx(filepath, "Hello DOCX")

    parser = DocxParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="sample.docx")

    assert isinstance(doc, Document)
    assert "Hello DOCX" in doc.text
    assert doc.metadata["parser"] == "DocxParser"


def test_docx_parser_empty(tmp_workspace):
    filepath = tmp_workspace / "empty.docx"
    _create_minimal_docx(filepath, "")

    parser = DocxParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="empty.docx")
    assert doc.text == ""


def test_docx_parser_corrupt(tmp_workspace):
    filepath = tmp_workspace / "corrupt.docx"
    filepath.write_bytes(b"not a docx file at all")

    parser = DocxParser()
    with pytest.raises(Exception):
        parser.parse(filepath.read_bytes(), filename=filepath.name)


# ------------------------------------------------------------------
# PptxParser
# ------------------------------------------------------------------

def _create_minimal_pptx(filepath: Path, texts: list[str] = None):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    if texts is None:
        texts = ["Slide 1 Title"]
    for slide_text in texts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        txBox.text_frame.text = slide_text
    prs.save(str(filepath))


def test_pptx_parser_valid(tmp_workspace):
    filepath = tmp_workspace / "sample.pptx"
    _create_minimal_pptx(filepath, ["Slide One", "Slide Two"])

    parser = PptxParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="sample.pptx")

    assert isinstance(doc, Document)
    assert "Slide One" in doc.text
    assert "Slide Two" in doc.text
    assert doc.page_count == 2
    assert doc.metadata["parser"] == "PptxParser"


def test_pptx_parser_empty(tmp_workspace):
    filepath = tmp_workspace / "empty.pptx"
    _create_minimal_pptx(filepath, [""])

    parser = PptxParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="empty.pptx")
    assert doc.text.strip() == ""


def test_pptx_parser_corrupt(tmp_workspace):
    filepath = tmp_workspace / "corrupt.pptx"
    filepath.write_bytes(b"definitely not pptx")

    parser = PptxParser()
    with pytest.raises(Exception):
        parser.parse(filepath.read_bytes(), filename=filepath.name)


# ------------------------------------------------------------------
# ImageParser (OCR via pytesseract)
# ------------------------------------------------------------------

@patch("parsers.image_parser.pytesseract.image_to_string")
def test_image_parser_valid(mock_ocr, tmp_workspace):
    """Create a simple white image and run OCR (expect empty or minimal text)."""
    from PIL import Image
    mock_ocr.return_value = "Mocked OCR text"

    filepath = tmp_workspace / "blank.png"
    img = Image.new("RGB", (100, 100), color="white")
    img.save(str(filepath))

    parser = ImageParser()
    doc = parser.parse(filepath.read_bytes(), filename=filepath.name, source_path="blank.png")

    assert isinstance(doc, Document)
    assert doc.metadata["parser"] == "ImageParser"
    assert doc.metadata["width"] == 100
    assert doc.metadata["height"] == 100


def test_image_parser_corrupt(tmp_workspace):
    filepath = tmp_workspace / "corrupt.png"
    filepath.write_bytes(b"not an image")

    parser = ImageParser()
    with pytest.raises(Exception):
        parser.parse(filepath.read_bytes(), filename=filepath.name)


def test_image_parser_caption_provider_extension():
    """Verify the CaptionProvider ABC cannot be instantiated."""
    with pytest.raises(TypeError):
        CaptionProvider()


# ------------------------------------------------------------------
# Factory
# ------------------------------------------------------------------

@pytest.mark.parametrize("ext, expected_parser_class", [
    (".txt", TextParser),
    (".md", TextParser),
    (".markdown", TextParser),
    (".pdf", PdfParser),
    (".docx", DocxParser),
    (".pptx", PptxParser),
    (".jpg", ImageParser),
    (".jpeg", ImageParser),
    (".png", ImageParser),
    (".gif", ImageParser),
    (".webp", ImageParser),
])
def test_factory_returns_correct_parser(ext, expected_parser_class):
    parser = get_parser(ext)
    assert isinstance(parser, expected_parser_class)


@pytest.mark.parametrize("ext", [".csv", ".exe", ".zip", ".mp3", ".html"])
def test_factory_rejects_unsupported(ext):
    with pytest.raises(ValueError, match="No parser registered"):
        get_parser(ext)
