"""
Parser for Microsoft PowerPoint files (.pptx) using ``python-pptx``.

Extracts text from every shape on every slide.
"""
from pathlib import Path
from typing import Optional

from pptx import Presentation

from parsers.base import BaseParser, Document
from utils.logger import logger


class PptxParser(BaseParser):
    """Handles .pptx files."""

    def parse(self, file_content: bytes, filename: str, source_path: Optional[str] = None) -> Document:
        """
        Extract text from each slide's text frames.

        Slide boundaries are separated by double newlines so that the
        downstream chunker can detect page/slide breaks.
        """
        import io
        prs = Presentation(io.BytesIO(file_content))
        slides_text: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

            if slide_parts:
                slides_text.append("\n".join(slide_parts))

        full_text = "\n\n".join(slides_text)
        slide_count = len(prs.slides)
        rel = source_path or filename

        meta: dict = {"parser": "PptxParser", "slide_count": slide_count}

        if not full_text.strip():
            logger.warning(f"PptxParser: '{rel}' contains no extractable text")

        logger.info(f"PptxParser: parsed '{rel}' ({slide_count} slides, {len(full_text)} chars)")

        return Document(
            filename=filename,
            source_path=rel,
            text=full_text,
            page_count=slide_count,
            metadata=meta,
        )
